"""Format Parser Lambda

Extracts text from documents (PDF, DOCX, PPTX, DWG, Markdown, TXT, Excel, CSV).
- PDF: Uses pypdf with graphics stripping for efficiency
- DOCX: Uses python-docx for text extraction
- PPTX: Uses python-pptx for text, LibreOffice for PDF conversion, pypdfium2 for images
- DWG: Uses aspose-cad for DWG-to-DXF conversion, ezdxf for text extraction + PNG rendering
- Excel/CSV: openpyxl (xlsx), xlrd (xls), csv module - per-sheet markdown tables
- Markdown/TXT: Direct text reading

Text files are chunked (4000 chars, 200 overlap) for optimal processing.
Saves result to S3 as format-parser/result.json for segment-builder to merge.
"""
import csv
import json
import os
import re
import subprocess
import tempfile

import pypdf
from pypdf.generic import DecodedStreamObject, NameObject

# Text chunking configuration
TEXT_CHUNK_SIZE = 15000
TEXT_CHUNK_OVERLAP = 500

# File type constants
TEXT_MIME_TYPES = (
    'text/plain',
    'text/markdown',
)

SPREADSHEET_MIME_TYPES = (
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
    'application/vnd.ms-excel',
    'text/csv',
)

PRESENTATION_MIME_TYPES = (
    'application/vnd.openxmlformats-officedocument.presentationml.presentation',
    'application/vnd.ms-powerpoint',
)

OFFICE_DOC_MIME_TYPES = (
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    'application/msword',
)

DXF_MIME_TYPES = (
    'application/dxf',
    'image/vnd.dxf',
)

from shared.ddb_client import (
    update_workflow_status,
    get_entity_prefix,
        WorkflowStatus,
    record_step_start,
    record_step_complete,
    record_step_error,
    record_step_skipped,
    StepName,
)
from shared.s3_analysis import get_s3_client, parse_s3_uri

BT_ET_PATTERN = re.compile(rb'BT\b.*?ET\b', re.DOTALL)


def get_document_base_path(file_uri: str) -> tuple:
    """Extract bucket and document base path from file URI."""
    bucket, key = parse_s3_uri(file_uri)
    key_parts = key.split('/')

    if 'documents' in key_parts:
        doc_idx = key_parts.index('documents')
        base_path = '/'.join(key_parts[:doc_idx + 2])
    else:
        base_path = '/'.join(key_parts[:-1])

    return bucket, base_path


def strip_graphics_inplace(reader: pypdf.PdfReader):
    """Strip non-text content from PDF in-place, keeping only BT..ET text blocks."""
    for page_num, page in enumerate(reader.pages):
        contents = page.get('/Contents')
        if contents is None:
            continue

        try:
            obj = contents.get_object()
            if isinstance(obj, pypdf.generic.ArrayObject):
                raw = b''.join(item.get_object().get_data() for item in obj)
                blocks = BT_ET_PATTERN.findall(raw)
                new_obj = DecodedStreamObject()
                new_obj.set_data(b'\n'.join(blocks))
                page[NameObject('/Contents')] = new_obj
            else:
                raw = obj.get_data()
                blocks = BT_ET_PATTERN.findall(raw)
                obj.set_data(b'\n'.join(blocks))
        except Exception as e:
            print(f'[format-parser] Page {page_num}: strip_graphics failed: {e}')


def process_pdf(file_uri: str) -> dict:
    """Download PDF, extract text per page, and save as result.json."""
    s3_client = get_s3_client()
    bucket, key = parse_s3_uri(file_uri)

    with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
        tmp_path = tmp.name
        s3_client.download_file(bucket, key, tmp_path)

    file_size = os.path.getsize(tmp_path)
    print(f'[format-parser] Downloaded PDF: {file_size} bytes')

    try:
        with open(tmp_path, 'rb') as f:
            reader = pypdf.PdfReader(f)
            num_pages = len(reader.pages)
            print(f'[format-parser] PDF has {num_pages} pages')

            print(f'[format-parser] Stripping graphics from content streams...')
            strip_graphics_inplace(reader)

            pages = []
            total_chars = 0

            print(f'[format-parser] Extracting text...')
            for page_num, page in enumerate(reader.pages):
                try:
                    text = (page.extract_text() or '').strip()
                except Exception as e:
                    print(f'[format-parser] Page {page_num}: extract_text failed: {e}')
                    text = ''
                pages.append({
                    'page_index': page_num,
                    'text': text,
                })
                total_chars += len(text)
                if (page_num + 1) % 500 == 0:
                    print(f'[format-parser] Extracted {page_num + 1}/{num_pages} pages')

        print(f'[format-parser] Done: {len(pages)} pages, {total_chars} chars')

        # Save to format-parser/result.json
        result_data = {'pages': pages}
        doc_bucket, base_path = get_document_base_path(file_uri)
        result_key = f'{base_path}/format-parser/result.json'

        s3_client.put_object(
            Bucket=doc_bucket,
            Key=result_key,
            Body=json.dumps(result_data, ensure_ascii=False),
            ContentType='application/json',
        )
        print(f'[format-parser] Saved result to s3://{doc_bucket}/{result_key}')

        return {
            'status': 'completed',
            'page_count': len(pages),
            'total_chars': total_chars,
        }

    finally:
        if os.path.exists(tmp_path):
            os.remove(tmp_path)


def chunk_text(text: str, chunk_size: int = TEXT_CHUNK_SIZE, overlap: int = TEXT_CHUNK_OVERLAP) -> list[str]:
    """Split text into overlapping chunks."""
    if not text:
        return []

    chunks = []
    start = 0
    text_len = len(text)

    while start < text_len:
        end = start + chunk_size

        # Try to break at sentence boundary
        if end < text_len:
            # Look for sentence end within last 200 chars
            search_start = max(start + chunk_size - 200, start)
            search_text = text[search_start:end]

            # Find last sentence boundary
            for sep in ['. ', '.\n', '? ', '?\n', '! ', '!\n', '\n\n']:
                last_sep = search_text.rfind(sep)
                if last_sep != -1:
                    end = search_start + last_sep + len(sep)
                    break

        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)

        # Move start with overlap
        start = end - overlap if end < text_len else text_len

    return chunks


def extract_docx_text(file_path: str) -> str:
    """Extract text from DOCX file using python-docx."""
    try:
        from docx import Document
        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return '\n\n'.join(paragraphs)
    except ImportError:
        print('python-docx not installed, falling back to basic extraction')
        return ''
    except Exception as e:
        print(f'Error extracting DOCX text: {e}')
        return ''


def process_text_file(file_uri: str, file_type: str) -> dict:
    """Process text-based file (DOCX, Markdown, TXT, CSV) and save chunks."""
    s3_client = get_s3_client()
    bucket, key = parse_s3_uri(file_uri)

    # Determine file extension
    ext = '.' + key.split('.')[-1].lower() if '.' in key else '.txt'

    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        tmp_path = tmp.name
        s3_client.download_file(bucket, key, tmp_path)

    file_size = os.path.getsize(tmp_path)
    print(f'[format-parser] Downloaded text file: {file_size} bytes, type: {file_type}')

    try:
        # Plain text, markdown, csv - direct read
        with open(tmp_path, 'r', encoding='utf-8') as f:
            text = f.read()

        if not text:
            print('[format-parser] No text extracted')
            chunks_data = [{'chunk_index': 0, 'text': ''}]
        else:
            # Chunk the text
            chunks = chunk_text(text)
            print(f'[format-parser] Text split into {len(chunks)} chunks')
            chunks_data = [{'chunk_index': i, 'text': chunk} for i, chunk in enumerate(chunks)]

        total_chars = sum(len(c['text']) for c in chunks_data)
        print(f'[format-parser] Done: {len(chunks_data)} chunks, {total_chars} chars')

        # Save to format-parser/result.json
        result_data = {'chunks': chunks_data}
        doc_bucket, base_path = get_document_base_path(file_uri)
        result_key = f'{base_path}/format-parser/result.json'

        s3_client.put_object(
            Bucket=doc_bucket,
            Key=result_key,
            Body=json.dumps(result_data, ensure_ascii=False),
            ContentType='application/json',
        )
        print(f'[format-parser] Saved result to s3://{doc_bucket}/{result_key}')

        return {
            'status': 'completed',
            'chunk_count': len(chunks_data),
            'total_chars': total_chars,
        }

    finally:
        if os.path.exists(tmp_path):
            os.remove(tmp_path)


def is_text_file(file_type: str) -> bool:
    """Check if file type is a text-based document."""
    return file_type in TEXT_MIME_TYPES


def is_spreadsheet(file_type: str) -> bool:
    """Check if file type is a spreadsheet."""
    return file_type in SPREADSHEET_MIME_TYPES


def is_presentation(file_type: str) -> bool:
    """Check if file type is a presentation."""
    return file_type in PRESENTATION_MIME_TYPES


def is_office_doc(file_type: str) -> bool:
    """Check if file type is a Word document."""
    return file_type in OFFICE_DOC_MIME_TYPES


def extract_pptx_text(file_path: str) -> list[dict]:
    """Extract text from PPTX file slide by slide using python-pptx.

    Returns list of dicts: [{'slide_index': 0, 'text': '...'}]
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation(file_path)
    slides = []

    for slide_idx, slide in enumerate(prs.slides):
        parts = []

        for shape in slide.shapes:
            # Text frames (titles, body text, text boxes)
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        parts.append(text)

            # Tables
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_texts = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_texts.append(cell_text)
                    if row_texts:
                        parts.append(' | '.join(row_texts))

        # Speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                parts.append(f'[Notes] {notes_text}')

        slide_text = '\n'.join(parts)
        slides.append({
            'slide_index': slide_idx,
            'text': slide_text,
        })

    return slides


def process_pptx(file_uri: str) -> dict:
    """Process PPTX: extract text, convert to PDF, render slide images."""
    import io
    import pypdfium2 as pdfium
    from PIL import Image

    s3_client = get_s3_client()
    bucket, key = parse_s3_uri(file_uri)

    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp_path = tmp.name
        s3_client.download_file(bucket, key, tmp_path)

    file_size = os.path.getsize(tmp_path)
    print(f'[format-parser] Downloaded PPTX: {file_size} bytes')

    tmp_dir = tempfile.mkdtemp()

    try:
        # 1. Extract text from PPTX using python-pptx
        slide_texts = extract_pptx_text(tmp_path)
        print(f'[format-parser] Extracted text from {len(slide_texts)} slides')

        # 2. Convert PPTX to PDF using LibreOffice
        result = subprocess.run(
            [
                'soffice', '--headless', '--norestore',
                '--convert-to', 'pdf',
                '--outdir', tmp_dir,
                tmp_path,
            ],
            capture_output=True,
            text=True,
            timeout=600,
        )
        if result.returncode != 0:
            print(f'[format-parser] LibreOffice stderr: {result.stderr}')
            raise RuntimeError(f'LibreOffice conversion failed: {result.stderr}')

        # Find the output PDF
        pdf_name = os.path.splitext(os.path.basename(tmp_path))[0] + '.pdf'
        pdf_path = os.path.join(tmp_dir, pdf_name)
        if not os.path.exists(pdf_path):
            # Fallback: look for any PDF in the output dir
            pdf_files = [f for f in os.listdir(tmp_dir) if f.endswith('.pdf')]
            if pdf_files:
                pdf_path = os.path.join(tmp_dir, pdf_files[0])
            else:
                raise RuntimeError('PDF output not found after LibreOffice conversion')

        print(f'[format-parser] Converted to PDF: {os.path.getsize(pdf_path)} bytes')

        # 3. Render PDF pages to PNG images using pypdfium2
        doc_bucket, base_path = get_document_base_path(file_uri)
        doc = pdfium.PdfDocument(pdf_path)
        page_count = len(doc)
        print(f'[format-parser] PDF has {page_count} pages')

        scale = 150 / 72  # 150 DPI
        pages = []
        total_chars = 0

        for page_idx in range(page_count):
            page = doc[page_idx]
            bitmap = page.render(scale=scale)
            img = bitmap.to_pil()

            # Convert to PNG bytes
            buf = io.BytesIO()
            img.save(buf, format='PNG')
            png_bytes = buf.getvalue()

            # Upload slide image to S3
            image_key = f'{base_path}/format-parser/slides/slide_{page_idx:04d}.png'
            s3_client.put_object(
                Bucket=doc_bucket,
                Key=image_key,
                Body=png_bytes,
                ContentType='image/png',
            )
            image_uri = f's3://{doc_bucket}/{image_key}'

            bitmap.close()
            page.close()

            # Match slide text by index
            text = ''
            if page_idx < len(slide_texts):
                text = slide_texts[page_idx].get('text', '')

            pages.append({
                'page_index': page_idx,
                'text': text,
                'image_uri': image_uri,
            })
            total_chars += len(text)

            if (page_idx + 1) % 50 == 0:
                print(f'[format-parser] Processed {page_idx + 1}/{page_count} slides')

        doc.close()

        print(f'[format-parser] Done: {len(pages)} slides, {total_chars} chars')

        # Save result.json
        result_data = {'pages': pages}
        result_key = f'{base_path}/format-parser/result.json'
        s3_client.put_object(
            Bucket=doc_bucket,
            Key=result_key,
            Body=json.dumps(result_data, ensure_ascii=False),
            ContentType='application/json',
        )
        print(f'[format-parser] Saved result to s3://{doc_bucket}/{result_key}')

        return {
            'status': 'completed',
            'page_count': len(pages),
            'total_chars': total_chars,
        }

    finally:
        if os.path.exists(tmp_path):
            os.remove(tmp_path)
        # Clean up temp dir
        import shutil
        if os.path.exists(tmp_dir):
            shutil.rmtree(tmp_dir, ignore_errors=True)


def process_docx(file_uri: str) -> dict:
    """Process DOCX: convert to PDF via LibreOffice, extract text + page images."""
    import io
    import pypdfium2 as pdfium

    s3_client = get_s3_client()
    bucket, key = parse_s3_uri(file_uri)

    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp_path = tmp.name
        s3_client.download_file(bucket, key, tmp_path)

    file_size = os.path.getsize(tmp_path)
    print(f'[format-parser] Downloaded DOCX: {file_size} bytes')

    tmp_dir = tempfile.mkdtemp()

    try:
        # 1. Convert DOCX to PDF using LibreOffice
        result = subprocess.run(
            [
                'soffice', '--headless', '--norestore',
                '--convert-to', 'pdf',
                '--outdir', tmp_dir,
                tmp_path,
            ],
            capture_output=True,
            text=True,
            timeout=600,
        )
        if result.returncode != 0:
            print(f'[format-parser] LibreOffice stderr: {result.stderr}')
            raise RuntimeError(f'LibreOffice conversion failed: {result.stderr}')

        # Find the output PDF
        pdf_name = os.path.splitext(os.path.basename(tmp_path))[0] + '.pdf'
        pdf_path = os.path.join(tmp_dir, pdf_name)
        if not os.path.exists(pdf_path):
            pdf_files = [f for f in os.listdir(tmp_dir) if f.endswith('.pdf')]
            if pdf_files:
                pdf_path = os.path.join(tmp_dir, pdf_files[0])
            else:
                raise RuntimeError('PDF output not found after LibreOffice conversion')

        print(f'[format-parser] Converted DOCX to PDF: {os.path.getsize(pdf_path)} bytes')

        # 2. Extract text per page using pypdf (with graphics stripping)
        with open(pdf_path, 'rb') as f:
            reader = pypdf.PdfReader(f)
            num_pages = len(reader.pages)
            print(f'[format-parser] PDF has {num_pages} pages')

            strip_graphics_inplace(reader)

            page_texts = []
            for page_num, page in enumerate(reader.pages):
                text = (page.extract_text() or '').strip()
                page_texts.append(text)

        # 3. Render PDF pages to PNG images using pypdfium2
        doc_bucket, base_path = get_document_base_path(file_uri)
        doc = pdfium.PdfDocument(pdf_path)
        page_count = len(doc)

        scale = 150 / 72  # 150 DPI
        pages = []
        total_chars = 0

        for page_idx in range(page_count):
            page = doc[page_idx]
            bitmap = page.render(scale=scale)
            img = bitmap.to_pil()

            buf = io.BytesIO()
            img.save(buf, format='PNG')
            png_bytes = buf.getvalue()

            image_key = f'{base_path}/format-parser/slides/slide_{page_idx:04d}.png'
            s3_client.put_object(
                Bucket=doc_bucket,
                Key=image_key,
                Body=png_bytes,
                ContentType='image/png',
            )
            image_uri = f's3://{doc_bucket}/{image_key}'

            bitmap.close()
            page.close()

            text = page_texts[page_idx] if page_idx < len(page_texts) else ''
            pages.append({
                'page_index': page_idx,
                'text': text,
                'image_uri': image_uri,
            })
            total_chars += len(text)

            if (page_idx + 1) % 50 == 0:
                print(f'[format-parser] Processed {page_idx + 1}/{page_count} pages')

        doc.close()

        print(f'[format-parser] Done: {len(pages)} pages, {total_chars} chars')

        # Save result.json
        result_data = {'pages': pages}
        result_key = f'{base_path}/format-parser/result.json'
        s3_client.put_object(
            Bucket=doc_bucket,
            Key=result_key,
            Body=json.dumps(result_data, ensure_ascii=False),
            ContentType='application/json',
        )
        print(f'[format-parser] Saved result to s3://{doc_bucket}/{result_key}')

        return {
            'status': 'completed',
            'page_count': len(pages),
            'total_chars': total_chars,
        }

    finally:
        if os.path.exists(tmp_path):
            os.remove(tmp_path)
        import shutil
        if os.path.exists(tmp_dir):
            shutil.rmtree(tmp_dir, ignore_errors=True)


def _sanitize_cell(value: str) -> str:
    """Sanitize cell value for markdown table (no newlines, escape pipes)."""
    return value.replace('\r\n', ' ').replace('\n', ' ').replace('\r', ' ').replace('|', '\\|').strip()


def _rows_to_markdown(rows: list[list[str]], sheet_name: str) -> str:
    """Convert rows (list of lists) to a markdown table with a sheet heading."""
    if not rows:
        return f'## Sheet: {sheet_name}\n\n(empty)'

    # Determine column count from the widest row
    col_count = max(len(r) for r in rows)

    # Pad rows to uniform width and sanitize cells
    padded = [[_sanitize_cell(c) for c in r] + [''] * (col_count - len(r)) for r in rows]

    header = padded[0]
    md_lines = [
        f'## Sheet: {sheet_name}',
        '',
        '| ' + ' | '.join(header) + ' |',
        '| ' + ' | '.join(['---'] * col_count) + ' |',
    ]
    for row in padded[1:]:
        md_lines.append('| ' + ' | '.join(row) + ' |')

    return '\n'.join(md_lines)


def _read_xlsx_sheets(file_path: str) -> list[tuple[str, list[list[str]]]]:
    """Read all sheets from .xlsx file using openpyxl. Returns [(sheet_name, rows)]."""
    import openpyxl
    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    sheets = []
    for name in wb.sheetnames:
        ws = wb[name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            str_row = [str(cell) if cell is not None else '' for cell in row]
            # Skip completely empty rows
            if any(c for c in str_row):
                rows.append(str_row)
        sheets.append((name, rows))
    wb.close()
    return sheets


def _read_xls_sheets(file_path: str) -> list[tuple[str, list[list[str]]]]:
    """Read all sheets from .xls file using xlrd. Returns [(sheet_name, rows)]."""
    import xlrd
    wb = xlrd.open_workbook(file_path)
    sheets = []
    for sheet in wb.sheets():
        rows = []
        for row_idx in range(sheet.nrows):
            str_row = [str(cell.value) if cell.value is not None else '' for cell in sheet.row(row_idx)]
            if any(c for c in str_row):
                rows.append(str_row)
        sheets.append((sheet.name, rows))
    return sheets


def _read_csv_sheets(file_path: str) -> list[tuple[str, list[list[str]]]]:
    """Read CSV file as a single sheet. Returns [(sheet_name, rows)]."""
    rows = []
    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
        reader = csv.reader(f)
        for row in reader:
            str_row = [cell.strip() for cell in row]
            if any(c for c in str_row):
                rows.append(str_row)
    sheet_name = os.path.splitext(os.path.basename(file_path))[0]
    return [(sheet_name, rows)]


def process_spreadsheet(file_uri: str, file_type: str) -> dict:
    """Process spreadsheet file (xlsx, xls, csv) and save per-sheet markdown tables."""
    s3_client = get_s3_client()
    bucket, key = parse_s3_uri(file_uri)

    ext = '.' + key.split('.')[-1].lower() if '.' in key else '.csv'

    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        tmp_path = tmp.name
        s3_client.download_file(bucket, key, tmp_path)

    file_size = os.path.getsize(tmp_path)
    print(f'[format-parser] Downloaded spreadsheet: {file_size} bytes, type: {file_type}')

    try:
        if ext == '.xlsx':
            sheets = _read_xlsx_sheets(tmp_path)
        elif ext == '.xls':
            sheets = _read_xls_sheets(tmp_path)
        else:
            sheets = _read_csv_sheets(tmp_path)

        print(f'[format-parser] Found {len(sheets)} sheet(s)')

        chunks_data = []
        for i, (sheet_name, rows) in enumerate(sheets):
            md_table = _rows_to_markdown(rows, sheet_name)
            chunks_data.append({'chunk_index': i, 'text': md_table})

        if not chunks_data:
            chunks_data = [{'chunk_index': 0, 'text': ''}]

        total_chars = sum(len(c['text']) for c in chunks_data)
        print(f'[format-parser] Done: {len(chunks_data)} sheets, {total_chars} chars')

        result_data = {'chunks': chunks_data}
        doc_bucket, base_path = get_document_base_path(file_uri)
        result_key = f'{base_path}/format-parser/result.json'

        s3_client.put_object(
            Bucket=doc_bucket,
            Key=result_key,
            Body=json.dumps(result_data, ensure_ascii=False),
            ContentType='application/json',
        )
        print(f'[format-parser] Saved result to s3://{doc_bucket}/{result_key}')

        return {
            'status': 'completed',
            'chunk_count': len(chunks_data),
            'total_chars': total_chars,
        }

    finally:
        if os.path.exists(tmp_path):
            os.remove(tmp_path)


def is_dxf(file_type: str) -> bool:
    """Check if file type is a DXF file."""
    return file_type in DXF_MIME_TYPES


def _extract_dxf_layout_text(layout) -> str:
    """Extract text entities from a single ezdxf layout."""
    parts = []
    for entity in layout:
        dxftype = entity.dxftype()
        if dxftype == 'TEXT':
            text = entity.dxf.text.strip()
            if text:
                parts.append(text)
        elif dxftype == 'MTEXT':
            text = entity.plain_text().strip()
            if text:
                parts.append(text)
        elif dxftype == 'ATTRIB':
            tag = entity.dxf.tag.strip()
            text = entity.dxf.text.strip()
            if text:
                parts.append(f'{tag}: {text}' if tag else text)
        elif dxftype == 'DIMENSION':
            text = getattr(entity.dxf, 'text', '').strip()
            if text:
                parts.append(text)
    return '\n'.join(parts)


def _extract_dxf_metadata(doc) -> str:
    """Extract metadata summary (layers, blocks) from ezdxf document."""
    lines = []

    layer_names = [layer.dxf.name for layer in doc.layers]
    if layer_names:
        lines.append(f'Layers ({len(layer_names)}): {", ".join(layer_names[:50])}')
        if len(layer_names) > 50:
            lines.append(f'  ... and {len(layer_names) - 50} more')

    block_names = [
        block.name for block in doc.blocks
        if not block.name.startswith('*')
    ]
    if block_names:
        lines.append(f'Blocks ({len(block_names)}): {", ".join(block_names[:50])}')
        if len(block_names) > 50:
            lines.append(f'  ... and {len(block_names) - 50} more')

    return '\n'.join(lines)


def process_dxf(file_uri: str) -> dict:
    """Process DXF: extract text + render PNG via ezdxf."""
    import shutil
    import ezdxf

    s3_client = get_s3_client()
    bucket, key = parse_s3_uri(file_uri)

    with tempfile.NamedTemporaryFile(suffix='.dxf', delete=False) as tmp:
        tmp_path = tmp.name
        s3_client.download_file(bucket, key, tmp_path)

    file_size = os.path.getsize(tmp_path)
    print(f'[format-parser] Downloaded DXF: {file_size} bytes')

    tmp_dir = tempfile.mkdtemp()

    try:
        doc = ezdxf.readfile(tmp_path)

        # Collect layout names (Model Space + Paper Space layouts)
        layout_names = ['Model']
        for name in doc.layout_names():
            if name != 'Model':
                layout_names.append(name)

        print(f'[format-parser] Found {len(layout_names)} layouts: {layout_names}')

        # 4. Extract metadata summary
        metadata_text = _extract_dxf_metadata(doc)

        # 5. Process each layout
        doc_bucket, base_path = get_document_base_path(file_uri)
        pages = []
        total_chars = 0

        page_idx = 0
        for layout_name in layout_names:
            # Extract text via ezdxf
            layout = doc.modelspace() if layout_name == 'Model' else doc.layout(layout_name)
            layout_text = _extract_dxf_layout_text(layout)

            # Skip empty non-Model layouts (common for unused Paper Space)
            entity_count = len(list(layout))
            if layout_name != 'Model' and entity_count == 0 and not layout_text:
                print(f'[format-parser] Skipping empty layout: {layout_name}')
                continue

            if page_idx == 0 and metadata_text:
                layout_text = f'{metadata_text}\n\n{layout_text}' if layout_text else metadata_text

            # Render to PNG via ezdxf + matplotlib
            try:
                import matplotlib
                matplotlib.use('Agg')
                from ezdxf.addons.drawing import RenderContext, Frontend
                from ezdxf.addons.drawing.config import Configuration, LineweightPolicy
                from ezdxf.addons.drawing.matplotlib import MatplotlibBackend
                import matplotlib.pyplot as plt

                fig = plt.figure(dpi=300)
                ax = fig.add_axes([0, 0, 1, 1])
                ax.set_aspect('equal')

                config = Configuration.defaults().with_changes(
                    lineweight_scaling=0.5,
                    min_lineweight=0.1,
                    lineweight_policy=LineweightPolicy.RELATIVE,
                )

                ctx = RenderContext(doc)
                out = MatplotlibBackend(ax)
                Frontend(ctx, out, config=config).draw_layout(layout, finalize=True)

                png_path = os.path.join(tmp_dir, f'layout_{page_idx:04d}.png')
                fig.savefig(png_path, dpi=300, bbox_inches='tight', pad_inches=0.1)
                plt.close(fig)

                with open(png_path, 'rb') as f:
                    png_bytes = f.read()

                image_key = f'{base_path}/format-parser/slides/layout_{page_idx:04d}.png'
                s3_client.put_object(
                    Bucket=doc_bucket,
                    Key=image_key,
                    Body=png_bytes,
                    ContentType='image/png',
                )
                image_uri = f's3://{doc_bucket}/{image_key}'
            except Exception as e:
                print(f'[format-parser] Failed to render layout "{layout_name}": {e}')
                image_uri = ''

            pages.append({
                'page_index': page_idx,
                'text': layout_text,
                'image_uri': image_uri,
            })
            total_chars += len(layout_text)
            page_idx += 1

        print(f'[format-parser] Done: {len(pages)} layouts, {total_chars} chars')

        # Save result.json
        result_data = {'pages': pages}
        result_key = f'{base_path}/format-parser/result.json'
        s3_client.put_object(
            Bucket=doc_bucket,
            Key=result_key,
            Body=json.dumps(result_data, ensure_ascii=False),
            ContentType='application/json',
        )
        print(f'[format-parser] Saved result to s3://{doc_bucket}/{result_key}')

        return {
            'status': 'completed',
            'page_count': len(pages),
            'total_chars': total_chars,
        }

    finally:
        if os.path.exists(tmp_path):
            os.remove(tmp_path)
        if os.path.exists(tmp_dir):
            shutil.rmtree(tmp_dir, ignore_errors=True)


def handler(event, context):
    print(f'Event: {json.dumps(event)}')

    workflow_id = event.get('workflow_id')
    document_id = event.get('document_id')
    file_uri = event.get('file_uri')
    file_type = event.get('file_type', '')

    is_pdf = file_type == 'application/pdf'
    is_text = is_text_file(file_type)
    is_pptx = is_presentation(file_type)
    is_docx = is_office_doc(file_type)
    is_spread = is_spreadsheet(file_type)
    is_dxf_file = is_dxf(file_type)

    if not is_pdf and not is_text and not is_pptx and not is_docx and not is_spread and not is_dxf_file:
        print(f'Skipping unsupported file: {file_type}')
        record_step_skipped(workflow_id, StepName.FORMAT_PARSER, f'File type {file_type} is not supported')
        return {
            **event,
            'format_parser': {
                'status': 'skipped',
                'reason': f'File type {file_type} is not supported'
            }
        }

    try:
        record_step_start(workflow_id, StepName.FORMAT_PARSER)

        if is_dxf_file:
            result = process_dxf(file_uri=file_uri)
        elif is_pptx:
            result = process_pptx(file_uri=file_uri)
        elif is_docx:
            result = process_docx(file_uri=file_uri)
        elif is_pdf:
            result = process_pdf(file_uri=file_uri)
        elif is_spread:
            result = process_spreadsheet(file_uri=file_uri, file_type=file_type)
        else:
            result = process_text_file(file_uri=file_uri, file_type=file_type)

        print(f'Format parser completed: {result}')
        record_step_complete(workflow_id, StepName.FORMAT_PARSER)

        return {
            **event,
            'format_parser': result
        }

    except Exception as e:
        error_msg = str(e)
        print(f'Error in format parser: {error_msg}')
        import traceback
        traceback.print_exc()
        entity_type = get_entity_prefix(file_type)
        update_workflow_status(document_id, workflow_id, WorkflowStatus.FAILED, entity_type=entity_type, error=error_msg)
        record_step_error(workflow_id, StepName.FORMAT_PARSER, error_msg)

        return {
            **event,
            'format_parser': {
                'status': 'failed',
                'error': error_msg
            }
        }
